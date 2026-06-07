import os
import zipfile
import traceback
from pathlib import Path
import fnmatch

import pdfplumber
import docx
import openpyxl
from pptx import Presentation
import chardet

MAX_FILE_SIZE = 50 * 1024 * 1024

def safe_extract_zip(zip_path, extract_to, log_callback=None):
    with zipfile.ZipFile(zip_path, 'r') as zf:
        for member in zf.namelist():
            member_path = Path(extract_to) / member
            if not member_path.resolve().is_relative_to(Path(extract_to).resolve()):
                raise Exception(f"Zip slip attempt: {member}")
        zf.extractall(extract_to)
        if log_callback:
            log_callback(f"Extracted {len(zf.namelist())} files from ZIP")

def should_skip_file(fp: Path, skip_exts: set, include_patterns: list, exclude_patterns: list) -> bool:
    if fp.stat().st_size > MAX_FILE_SIZE:
        return True
    if fp.suffix.lower() in skip_exts:
        return True
    if include_patterns:
        if not any(fnmatch.fnmatch(fp.name, p) for p in include_patterns):
            return True
    if exclude_patterns:
        if any(fnmatch.fnmatch(fp.name, p) for p in exclude_patterns):
            return True
    return False

def detect_encoding(file_path: Path) -> str:
    with open(file_path, 'rb') as f:
        raw = f.read(min(1024*1024, os.path.getsize(file_path)))
        result = chardet.detect(raw)
        return result.get('encoding', 'utf-8')

def extract_text(fp: Path, log_callback=None) -> str:
    ext = fp.suffix.lower()
    try:
        if ext == '.pdf':
            with pdfplumber.open(fp) as pdf:
                text = "\n".join(p.extract_text() or "" for p in pdf.pages)
        elif ext == '.docx':
            doc = docx.Document(fp)
            text = "\n".join(p.text for p in doc.paragraphs)
        elif ext == '.xlsx':
            wb = openpyxl.load_workbook(fp, data_only=True)
            lines = []
            for sheet in wb.worksheets:
                lines.append(f"[Sheet: {sheet.title}]")
                for row in sheet.iter_rows(values_only=True):
                    lines.append("\t".join(str(c) for c in row if c))
            text = "\n".join(lines)
        elif ext == '.pptx':
            prs = Presentation(fp)
            lines = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        lines.append(shape.text)
            text = "\n".join(lines)
        else:
            encoding = detect_encoding(fp)
            with open(fp, 'r', encoding=encoding, errors='replace') as f:
                text = f.read()
        # Basic cleanup: replace null bytes and other nasties
        text = text.replace('\x00', '')
        return text.encode('utf-8', errors='replace').decode('utf-8')
    except Exception as e:
        err_msg = f"[ERROR extracting {fp.name}: {traceback.format_exc()}]"
        if log_callback:
            log_callback(err_msg)
        return err_msg.encode('utf-8', errors='replace').decode('utf-8')